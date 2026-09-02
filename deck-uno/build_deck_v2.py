#!/usr/bin/env python3
"""deck-uno v2: plan-first deck builder.

流程学习 dsh-univer-office：
  1) PLAN   —— 写 pptx 之前先规划整副 deck 的版式：
               封面(标题+副标) / 导览 / 章节页(复杂 deck 自动插入) /
               内容页结构轮换(卡片·对照·大数字+图表·流程·引言) / Q&A / 致谢
  2) SVG    —— 把版式模型渲染成逐页 SVG 线框草稿（排版预览 + 越界检查）
  3) PPTX   —— 同一个版式模型渲染为可编辑 pptx（图表用原生 chart）
"""
import json, os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- constants
NAVY, BLUE, TEAL, AMBER = "16294D", "2F6BFF", "12B5A5", "F2A93B"
BG, DARK, BODY, MUTED   = "FAFBFD", "1B2A41", "44546A", "7A8AA3"
LINE, WHITE, SOFT       = "D9E1EF", "FFFFFF", "9FB4D8"
SOFTBG                  = "EEF2F9"
ACCENTS                 = [BLUE, TEAL, AMBER]
W_EMU, H_EMU            = 12192000, 6858000
FONT                    = "Hiragino Sans GB"  # 本机 LibreOffice 可渲染的中文黑体
EMU_PER_PX              = 12700          # svg 草稿 960x540

# ================================================================ 1. PLAN
# 内容素材（提炼后：每条 = 「加粗观点：解释」，而非平铺段落）
CONTENT = [
 ("为什么不能只保存对话", "开场", "cards", [
   ("Context Length 只解决了“读懂当前”", "它的增长从未触碰“记住过去”"),
   ("保存对话本身 = 保存噪声", "过程、寒暄、已被推翻的尝试全被存档"),
   ("正确目标：保存被整理出来的知识", "而不是保存对话的原始流"),
   ("线性记录的三个痛点", "体积膨胀 · 检索困难 · 无法关联")]),
 ("什么是一个 Session？", "Session切分", "compare", [
   ("粗看不成问题", "一次启动到一次退出，就是一段对话"),
   ("现实更复杂", "一次会话换了三个不相干的话题；一个任务横跨多次重启"),
   ("单位若错了", "后面所有摘要和抽取都对着错误的对象工作"),
   ("结论", "框架给出的物理边界 ≠ 语义上的 session")]),
 ("切分靠语义信号，而非物理边界", "Session切分", "strength", [
   ("话题切换（最强）", "从 sqlite 存储跳到集群运维，就该切成两段"),
   ("任务边界（次强）", "一次明确交付的完成 = 一段叙事的自然终结"),
   ("时间间隔（最弱）", "中断一夜后的继续，值得检查是否切开"),
   ("物理边界仅作提示", "收敛则顺势切分，未收敛建议并入下一段")]),
 ("一个 session，一个主题", "Session切分", "quote", [
   ("准则", "一个 session 恰好容纳一个可独立陈述的主题"),
   ("太小", "摘要被迫丢弃异质信息"),
   ("太大", "窗口没有空间等待话题收敛"),
   ("切错可修复", "低层原文总在，更高层归档时可重新缝合；不识别才是不可修复")]),
 ("对数空间的 Session Summary", "分层归档", "chart", [
   ("若摘要和原对话一样长，摘要毫无意义", ""),
   ("目标：log₁₆(n) 的空间代价", "每上升一层压缩 16 倍"),
   ("被压缩的是叙事和过程", "不是结论")]),
 ("分层归档 = B 树逻辑", "分层归档", "compare", [
   ("拒绝一次性压缩", "“读完一万字写出两百字”，损失不可控"),
   ("分层归档", "每层只面对上一层的产物，输入已是高密度内容"),
   ("损失发生在每层边缘", "而非一次性的坍塌"),
   ("B 树类比", "不指望根节点存住所有数据，只指望从根出发总能找到")]),
 ("三层归档与下钻回溯", "分层归档", "flow", [
   ("L2 主题摘要", "最稀疏，面向主题检索"),
   ("L1 会话摘要", "16:1 归档"),
   ("L0 原始对话", "冷存，不删除，退出热路径"),
   ("精确回溯", "从高层逐层下钻，每下一层多一次检索")]),
 ("滑动窗口：抽取的时机", "知识抽取", "flow", [
   ("已归档", "窗口滑过 = 既成事实"),
   ("窗口内", "待整理的候选"),
   ("未来", "尚未发生的对话"),
   ("窗口的意义在于节奏", "而非容量")]),
 ("窗口大小 = 话题的生命周期", "知识抽取", "compare", [
   ("太短", "把猜测和已被推翻的尝试误当成知识写入"),
   ("太长", "整理的时机姗姗来迟"),
   ("实践", "窗口覆盖一个完整话题的自然生命周期"),
   ("梳理时机", "出现收敛迹象：结论被确认、话题被切换")]),
 ("知识点：脱离上下文仍成立", "知识抽取", "goodbad", [
   ("✗ 用户说要用 sqlite 换掉 duckdb", "只在语境里有意义——那是聊天记录"),
   ("✓ 存储层评估过 sqlite 单库方案，因 json contains 查询不可行而搁置", "任何时候可检索、可直接使用"),
   ("每次梳理都要裁决", "新增、修正还是否定；否定留下冲突痕迹，绝不悄悄覆盖")]),
 ("提取原则（一）：记什么", "提取原则", "cards", [
   ("记教训而非流水账", "“用 Arc<Mutex<T>>”合格；“用合适的同步机制”不合格"),
   ("对质量保守，对时机激进", "拿不准宁可跳过，但检查频率不能省"),
   ("纠错链是金子", "问题→错答→纠正→修正，最难从文档重新获得"),
   ("分类决定模板", "答对/纠错链/探索/bug/决策各有骨架")]),
 ("提取原则（二）：怎么记", "提取原则", "flow", [
   ("相对时间 → 绝对时间", "“明天提醒”→“2026-08-29 提醒”"),
   ("写给未来的读取者", "“明天”在入库那一刻就开始失效"),
   ("宁冗勿漏关系", "谓词不完美也要把三元组挂进图"),
   ("绝不存敏感信息", "密码、密钥、token 不进入任何一层")]),
 ("三元组：把隐含关系显式化", "知识图", "compare", [
   ("关系藏在语气和语序中", "“sqlite 搁置”与“json contains 是核心依赖”有因果"),
   ("整理时显式化", "每个知识点拆为 (主语, 谓语, 宾语) 写入 RDF 风格的图"),
   ("分工", "正文负责“读到时理解”，三元组负责“需要时找到”"),
   ("图的遍历", "是发现隐性关联的唯一手段")]),
 ("图的价值：可遍历性", "知识图", "quote", [
   ("承认噪声", "三元组抽取不完美，谓词粒度难以精确"),
   ("错误可被发现", "一条错误三元组是可挑战、可 resolve 的对象"),
   ("未结构化的对话", "连犯错的资格都没有")]),
 ("自动回忆：按 Agent 机制适配", "读取通道", "columns3", [
   ("每轮重构 system prompt", "确定性注入：任务开始执行一次 JSE 查询"),
   ("持续上下文、能自主调工具", "技能层检索：需要时自己发起，中途可补查"),
   ("被动响应、几乎无主动权", "边界触发：新任务、用户提醒时外部注入"),
   ("hypatia skill 的使命", "同一套查询能力，包装成每个 Agent 最顺手的形态")]),
 ("显式写入与显式读取", "读取通道", "compare", [
   ("显式写入", "“把运维文档写进 hypatia”——现成材料跳过窗口，直接拆解入库"),
   ("显式读取", "“读 hypatia 找运维记忆”——检索从 Agent 判断变成用户指令"),
   ("手动复位", "自动回忆失灵时的一句话复位：成本一句话，回报是整个知识库"),
   ("例", "opencode 要直接 ssh 服务器时，一句提醒就够了")]),
 ("全景：从线性对话流到知识图", "总结", "flow", [
   ("写入路径", "线性对话流 → 语义切分 → L0/L1/L2 对数归档"),
   ("抽取路径", "对话流经滑动窗口 → 知识点 + 三元组 → 知识图"),
   ("例外路径", "显式写入直通知识图，跳过窗口"),
   ("读取通道", "自动回忆（按 Agent 适配）+ 显式读取 → 服务 Agent")]),
]

SECTIONS = [  # 复杂 deck 的章节页（kicker → 章节定义；None 标题 = 并入上一章）
    ("Session切分", "01", "Session 的识别与切分", "整理之前，先回答“什么是一个 session”"),
    ("分层归档",   "02", "对数空间的 Session Summary", "如果一个摘要和原对话一样长，它就毫无意义"),
    ("知识抽取",   "03", "滑动窗口与知识点梳理", "压缩之外，更关键的是抽取"),
    ("提取原则",   "03", None, None),          # 并入 03
    ("知识图",     "04", "三元组与图结构", "知识点彼此不是孤立的"),
    ("读取通道",   "05", "自动回忆与显式读写", "整理是写的一侧，回忆是读的一侧"),
    ("总结",       "05", None, None),          # 并入 05
]

def plan_deck():
    """返回整副 deck 的版式计划（写任何 pptx 之前先确定）。"""
    plan = [
        {"layout": "cover", "page": 1, "title": "非线性的记忆",
         "subtitle": "从线性对话流到可检索的知识结构",
         "kicker": "HYPATIA · MEMORY DESIGN", "meta": "基于 docs/memory-nolinear.md · 50 min"},
        {"layout": "toc", "page": 2, "title": "今天的路线", "kicker": "OVERVIEW"},
    ]
    page, seen = 3, set()
    for (title, kicker, layout, points) in CONTENT:
        sec = next((s for s in SECTIONS if s[0] == kicker), None)
        if sec and sec[2] and sec[1] not in seen:      # 章节页
            seen.add(sec[1])
            plan.append({"layout": "section", "page": page, "num": sec[1],
                         "title": sec[2], "subtitle": sec[3]})
            page += 1
        plan.append({"layout": layout, "page": page, "title": title,
                     "kicker": kicker, "points": points})
        page += 1
    plan += [
        {"layout": "summary", "page": page, "title": "对话是给人看的，知识是给机器用的",
         "pills": [("对数分层", "控制体积"), ("滑动窗口", "控制时机"),
                   ("三元组图", "提供发现路径"), ("双通道", "覆盖日常与例外")],
         "note": "四件事拼在一起 = “有限记忆的 AI 使用无限长度的知识内容”"},
        {"layout": "qa", "page": page + 1, "title": "Q & A",
         "subtitle": "欢迎提问：切分 · 归档 · 抽取 · 图结构 · 读写通道"},
        {"layout": "thanks", "page": page + 2, "title": "谢谢",
         "subtitle": "非线性的记忆 · hypatia",
         "meta": "github.com/MarchLiu/hypatia"},
    ]
    total = plan[-1]["page"]
    for s in plan:
        s["total"] = total
    return plan

# ============================================== 2. 版式模型 → draw list
def px(v):  # emu → svg px
    return round(v / EMU_PER_PX, 1)

def draw_cover(s):
    return [
        {"k": "rect", "x": 0, "y": 0, "cx": W_EMU, "cy": H_EMU, "fill": NAVY},
        {"k": "oval", "x": 7683500, "y": 1270000, "cx": 4191000, "cy": 4191000, "stroke": BLUE},
        {"k": "oval", "x": 8280400, "y": 1866900, "cx": 2997200, "cy": 2997200, "stroke": TEAL},
        {"k": "oval", "x": 8864600, "y": 2451100, "cx": 1828800, "cy": 1828800, "stroke": AMBER},
        {"k": "rect", "x": 812800, "y": 711200, "cx": 330200, "cy": 50800, "fill": BLUE},
        {"k": "text", "x": 1231900, "y": 635000, "cx": 6000000, "cy": 300000,
         "runs": [[["HYPATIA · MEMORY DESIGN", 15, SOFT, True]]]},
        {"k": "text", "x": 812800, "y": 2286000, "cx": 6800000, "cy": 1800000,
         "runs": [[["非线性的记忆", 64, WHITE, True]]]},
        {"k": "text", "x": 812800, "y": 4000500, "cx": 6800000, "cy": 600000,
         "runs": [[["从线性对话流到可检索的知识结构", 24, SOFT, False]]]},
        {"k": "text", "x": 812800, "y": 5900000, "cx": 6400000, "cy": 350000,
         "runs": [[["基于 docs/memory-nolinear.md · 50 min", 15, MUTED, False]]]}]

def draw_header(title, kicker, page, total):
    return [{"k": "rect", "x": 0, "y": 0, "cx": W_EMU, "cy": H_EMU, "fill": BG},
            {"k": "rect", "x": 812800, "y": 660400, "cx": 330200, "cy": 50800, "fill": BLUE},
            {"k": "text", "x": 990600, "y": 600000, "cx": 6000000, "cy": 280000,
             "runs": [[[kicker, 13, BLUE, True]]]},
            {"k": "text", "x": 812800, "y": 1000000, "cx": 9500000, "cy": 700000,
             "runs": [[[title, 32, DARK, True]]]},
            {"k": "text", "x": 10591800, "y": 660400, "cx": 787400, "cy": 280000,
             "runs": [[[f"{page:02d} / {total}", 13, MUTED, False]]], "align": "right"},
            {"k": "rect", "x": 812800, "y": 1780000, "cx": 10566400, "cy": 12700, "fill": LINE}]

def draw_card(x, y, cx, cy, accent):
    return [{"k": "rrect", "x": x, "y": y, "cx": cx, "cy": cy, "fill": WHITE, "stroke": LINE},
            {"k": "rect", "x": x, "y": y + 60000, "cx": 60000, "cy": cy - 120000, "fill": accent}]

def draw_toc(s):
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    items = [
        ("01", "Session 的识别与切分", "语义信号：话题 · 任务 · 时间"),
        ("02", "对数空间的 Session Summary", "log₁₆(n) · 分层归档 · 下钻回溯"),
        ("03", "滑动窗口与知识点梳理", "时机 · 提取原则 · 纠错链"),
        ("04", "三元组与图结构", "显式化关系 · 图遍历发现"),
        ("05", "读写通道", "自动回忆适配 · 显式写入/读取")]
    top, ch, gap = 1905000, 790000, 120000
    y = top
    for i, (num, t, sub) in enumerate(items):
        acc = ACCENTS[i % 3]
        d += draw_card(812800, y, 10566400, ch, acc)
        d.append({"k": "text", "x": 1117600, "y": y + 150000, "cx": 900000, "cy": 450000,
                  "runs": [[[num, 28, acc, True]]]})
        d.append({"k": "text", "x": 2100000, "y": y + 110000, "cx": 8800000, "cy": 350000,
                  "runs": [[[t, 21, DARK, True]]]})
        d.append({"k": "text", "x": 2100000, "y": y + 460000, "cx": 8800000, "cy": 300000,
                  "runs": [[[sub, 15, MUTED, False]]]})
        y += ch + gap
    return d

def draw_section(s):
    return [
        {"k": "rect", "x": 0, "y": 0, "cx": W_EMU, "cy": H_EMU, "fill": NAVY},
        {"k": "text", "x": 700000, "y": 900000, "cx": 3000000, "cy": 2200000,
         "runs": [[[s["num"], 96, BLUE, True]]], "alpha": 0.35},
        {"k": "rect", "x": 812800, "y": 3200000, "cx": 600000, "cy": 60000, "fill": AMBER},
        {"k": "text", "x": 812800, "y": 3500000, "cx": 8000000, "cy": 800000,
         "runs": [[[s["title"], 40, WHITE, True]]]},
        {"k": "text", "x": 812800, "y": 4400000, "cx": 8000000, "cy": 500000,
         "runs": [[[s["subtitle"], 20, SOFT, False]]]}]

def draw_cards(s):                       # 观点卡片（加粗观点 + 弱化解释）
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    pts = s["points"]; n = len(pts)
    top, bottom, gap = 2032000, 6550000, 150000
    ch = (bottom - top - gap * (n - 1)) // n
    size = 20 if n <= 4 else 18
    y = top
    for i, (lead, rest) in enumerate(pts):
        acc = ACCENTS[i % 3]
        d += draw_card(812800, y, 10566400, ch, acc)
        runs = [[lead, size, DARK, True]]
        if rest:
            runs.append(["　" + rest, size - 3, BODY, False])
        d.append({"k": "text", "x": 1231900, "y": y, "cx": 9850000, "cy": ch,
                  "runs": [runs], "anchor": "middle"})
        y += ch + gap
    return d

def draw_compare(s):                     # 2x2 对照 + 底部结论条
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    pts = s["points"]
    cw, chh, gx, gy = 5143000, 1900000, 280000, 220000
    x0, y0 = 812800, 2032000
    for i, (lead, rest) in enumerate(pts[:4]):
        r, c = divmod(i, 2)
        acc = ACCENTS[i % 3]
        x, y = x0 + c * (cw + gx), y0 + r * (chh + gy)
        d += draw_card(x, y, cw, chh, acc)
        d.append({"k": "text", "x": x + 220000, "y": y + 160000, "cx": cw - 440000, "cy": 500000,
                  "runs": [[[lead, 19, DARK, True]]]})
        if rest:
            d.append({"k": "text", "x": x + 220000, "y": y + 720000, "cx": cw - 440000, "cy": chh - 880000,
                      "runs": [[[rest, 15, BODY, False]]]})
    if len(pts) > 4:
        yb = y0 + 2 * (chh + gy) + 60000
        d.append({"k": "rrect", "x": x0, "y": yb, "cx": 2 * cw + gx, "cy": 700000,
                  "fill": SOFTBG, "stroke": LINE})
        lead, rest = pts[4]
        d.append({"k": "text", "x": x0 + 300000, "y": yb, "cx": 2 * cw + gx - 600000, "cy": 700000,
                  "runs": [[[lead + ("：" + rest if rest else ""), 17, DARK, True]]],
                  "align": "center", "anchor": "middle"})
    return d

def draw_strength(s):                    # 强/中/弱三层横卡
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    colors, tags = [BLUE, TEAL, AMBER], ["强", "次", "弱"]
    top, ch, gap = 2032000, 1180000, 200000
    y = top
    for i, (lead, rest) in enumerate(s["points"][:3]):
        acc = colors[i]
        d += draw_card(812800, y, 10566400, ch, acc)
        d.append({"k": "rrect", "x": 1050000, "y": y + 320000, "cx": 800000, "cy": 540000,
                  "fill": acc})
        d.append({"k": "text", "x": 1050000, "y": y + 320000, "cx": 800000, "cy": 540000,
                  "runs": [[[tags[i], 18, WHITE, True]]], "align": "center", "anchor": "middle"})
        d.append({"k": "text", "x": 2100000, "y": y + 180000, "cx": 8600000, "cy": 450000,
                  "runs": [[[lead, 20, DARK, True]]]})
        d.append({"k": "text", "x": 2100000, "y": y + 640000, "cx": 8600000, "cy": 400000,
                  "runs": [[[rest, 15, BODY, False]]]})
        y += ch + gap
    y += 60000
    lead, rest = s["points"][3]
    d.append({"k": "text", "x": 812800, "y": y, "cx": 10566400, "cy": 500000,
              "runs": [[[lead + "：" + rest, 16, MUTED, False]]], "anchor": "middle"})
    return d

def draw_quote(s):                       # 深底观点页
    d = [{"k": "rect", "x": 0, "y": 0, "cx": W_EMU, "cy": H_EMU, "fill": NAVY},
         {"k": "rect", "x": 812800, "y": 711200, "cx": 330200, "cy": 50800, "fill": AMBER},
         {"k": "text", "x": 1231900, "y": 635000, "cx": 6000000, "cy": 300000,
          "runs": [[[s["kicker"], 15, SOFT, True]]]},
         {"k": "text", "x": 1117600, "y": 2050000, "cx": 10000000, "cy": 1300000,
          "runs": [[[s["points"][0][1] or s["points"][0][0], 34, WHITE, True]]], "align": "center"}]
    half = s["points"][1:3]
    cw = 4900000
    x = 1117600
    for lead, rest in half:
        d.append({"k": "rrect", "x": x, "y": 3700000, "cx": cw, "cy": 1000000,
                  "fill": "22355C", "stroke": "33507F"})
        d.append({"k": "text", "x": x + 200000, "y": 3860000, "cx": cw - 400000, "cy": 400000,
                  "runs": [[[lead, 18, WHITE, True]]], "align": "center"})
        if rest:
            d.append({"k": "text", "x": x + 200000, "y": 4280000, "cx": cw - 400000, "cy": 400000,
                      "runs": [[[rest, 14, SOFT, False]]], "align": "center"})
        x += cw + 350000
    if len(s["points"]) > 3:
        d.append({"k": "text", "x": 1117600, "y": 5300000, "cx": 9956400, "cy": 500000,
                  "runs": [[[s["points"][3][0] + "：" + s["points"][3][1], 15, MUTED, False]]],
                  "align": "center"})
    return d

def draw_chart(s):                       # 大数字 + 原生柱状图（对数压缩）
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    d.append({"k": "text", "x": 812800, "y": 2300000, "cx": 3500000, "cy": 1300000,
              "runs": [[["16 : 1", 72, BLUE, True]]]})
    d.append({"k": "text", "x": 812800, "y": 3700000, "cx": 3500000, "cy": 900000,
              "runs": [[["每上升一层\n压缩十六倍", 18, BODY, False]]]})
    d.append({"k": "chart", "x": 4700000, "y": 2050000, "cx": 6600000, "cy": 3900000,
              "title": "相对体积（对数压缩）",
              "cats": ["L2 主题摘要", "L1 会话摘要", "L0 原始对话"],
              "series": [("相对体积 %", [0.39, 6.25, 100.0])]})
    p = s["points"][1]
    d.append({"k": "text", "x": 812800, "y": 6150000, "cx": 10566400, "cy": 450000,
              "runs": [[[p[0] + ("：" + p[1] if p[1] else ""), 16, DARK, True]]],
              "anchor": "middle"})
    return d

def draw_flow(s):                        # 横向流程节点 + 箭头
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    pts = s["points"]; n = min(len(pts), 4)
    y0 = 2800000
    node_w, node_h = 2350000, 1500000
    gapx = (10566400 - n * node_w) // max(n - 1, 1) if n > 1 else 0
    x = 812800
    for i in range(n):
        lead, rest = pts[i]
        acc = ACCENTS[i % 3]
        d.append({"k": "rrect", "x": x, "y": y0, "cx": node_w, "cy": node_h,
                  "fill": WHITE, "stroke": acc, "sw": 28575})
        d.append({"k": "text", "x": x + 140000, "y": y0 + 200000, "cx": node_w - 280000, "cy": 450000,
                  "runs": [[[lead, 17, DARK, True]]], "align": "center"})
        if rest:
            d.append({"k": "text", "x": x + 140000, "y": y0 + 680000, "cx": node_w - 280000, "cy": 700000,
                      "runs": [[[rest, 13, BODY, False]]], "align": "center"})
        if i < n - 1:
            d.append({"k": "line", "x": x + node_w + 60000, "y": y0 + node_h // 2,
                      "cx": gapx - 120000, "cy": 0, "stroke": MUTED, "arrow": True})
        x += node_w + gapx
    if len(pts) > n:
        lead, rest = pts[n]
        d.append({"k": "text", "x": 812800, "y": 4800000, "cx": 10566400, "cy": 500000,
                  "runs": [[[lead + "：" + rest, 17, DARK, True]]], "anchor": "middle"})
    return d

def draw_goodbad(s):                     # ✗/✓ 对照
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    specs = [(MUTED, "✗"), (TEAL, "✓")]
    cw, chh = 5143000, 2500000
    for i in range(2):
        lead, rest = s["points"][i]
        border, mark = specs[i]
        x = 812800 + i * (cw + 280000)
        d.append({"k": "rrect", "x": x, "y": 2032000, "cx": cw, "cy": chh,
                  "fill": WHITE, "stroke": border, "sw": 28575})
        d.append({"k": "text", "x": x + 200000, "y": 2232000, "cx": 600000, "cy": 500000,
                  "runs": [[[mark, 28, border, True]]]})
        d.append({"k": "text", "x": x + 800000, "y": 2232000, "cx": cw - 1000000, "cy": 1000000,
                  "runs": [[[lead[2:].strip(), 17, DARK, True]]]})
        d.append({"k": "text", "x": x + 200000, "y": 3400000, "cx": cw - 400000, "cy": 900000,
                  "runs": [[[rest, 14, BODY, False]]]})
    lead, rest = s["points"][2]
    d.append({"k": "rrect", "x": 812800, "y": 4900000, "cx": 10566400, "cy": 800000,
              "fill": SOFTBG, "stroke": LINE})
    d.append({"k": "text", "x": 1117600, "y": 4900000, "cx": 9956400, "cy": 800000,
              "runs": [[[lead + "：" + rest, 17, DARK, True]]], "anchor": "middle"})
    return d

def draw_columns3(s):                    # 三竖卡
    d = draw_header(s["title"], s["kicker"], s["page"], s["total"])
    cw, chh, gap = 3380000, 3200000, 213000
    x = 812800
    for i, (lead, rest) in enumerate(s["points"][:3]):
        acc = ACCENTS[i % 3]
        d += draw_card(x, 2032000, cw, chh, acc)
        d.append({"k": "text", "x": x + 220000, "y": 2260000, "cx": cw - 440000, "cy": 1100000,
                  "runs": [[[lead, 18, DARK, True]]]})
        d.append({"k": "text", "x": x + 220000, "y": 3450000, "cx": cw - 440000, "cy": 1600000,
                  "runs": [[[rest, 15, BODY, False]]]})
        x += cw + gap
    lead, rest = s["points"][3]
    d.append({"k": "text", "x": 812800, "y": 5600000, "cx": 10566400, "cy": 500000,
              "runs": [[[lead + "：" + rest, 16, MUTED, False]]], "anchor": "middle"})
    return d

def _dark_hero(art=BLUE):
    return [{"k": "rect", "x": 0, "y": 0, "cx": W_EMU, "cy": H_EMU, "fill": NAVY},
            {"k": "oval", "x": 7683500, "y": 1270000, "cx": 4191000, "cy": 4191000, "stroke": art},
            {"k": "oval", "x": 8280400, "y": 1866900, "cx": 2997200, "cy": 2997200, "stroke": TEAL},
            {"k": "rect", "x": 812800, "y": 711200, "cx": 330200, "cy": 50800, "fill": AMBER}]

def draw_summary(s):
    d = _dark_hero(BLUE)
    d += [{"k": "text", "x": 1231900, "y": 635000, "cx": 6000000, "cy": 300000,
           "runs": [[["写在最后", 15, SOFT, True]]]},
          {"k": "text", "x": 812800, "y": 1900000, "cx": 10000000, "cy": 900000,
           "runs": [[[s["title"], 40, WHITE, True]]]}]
    x = 812800
    for a, b in s["pills"]:
        d.append({"k": "rrect", "x": x, "y": 3300000, "cx": 2500000, "cy": 900000,
                  "stroke": BLUE, "sw": 19050})
        d.append({"k": "text", "x": x, "y": 3300000, "cx": 2500000, "cy": 900000,
                  "runs": [[[a, 18, WHITE, True]], [[b, 13, SOFT, False]]],
                  "align": "center", "anchor": "middle"})
        x += 2700000
    d.append({"k": "text", "x": 812800, "y": 4600000, "cx": 10000000, "cy": 500000,
              "runs": [[[s["note"], 20, SOFT, False]]]})
    return d

def draw_qa(s):
    d = _dark_hero(TEAL)
    d += [{"k": "text", "x": 1231900, "y": 635000, "cx": 6000000, "cy": 300000,
           "runs": [[["DISCUSSION", 15, SOFT, True]]]},
          {"k": "text", "x": 812800, "y": 2350000, "cx": 8000000, "cy": 1200000,
           "runs": [[[s["title"], 60, WHITE, True]]]},
          {"k": "text", "x": 812800, "y": 3800000, "cx": 8000000, "cy": 500000,
           "runs": [[[s["subtitle"], 20, SOFT, False]]]}]
    return d

def draw_thanks(s):
    d = _dark_hero(AMBER)
    d += [{"k": "text", "x": 1231900, "y": 635000, "cx": 6000000, "cy": 300000,
           "runs": [[["THANK YOU", 15, SOFT, True]]]},
          {"k": "text", "x": 812800, "y": 2450000, "cx": 8000000, "cy": 1100000,
           "runs": [[[s["title"], 54, WHITE, True]]]},
          {"k": "text", "x": 812800, "y": 3800000, "cx": 8000000, "cy": 500000,
           "runs": [[[s["subtitle"], 20, SOFT, False]]]},
          {"k": "text", "x": 812800, "y": 5900000, "cx": 8000000, "cy": 350000,
           "runs": [[[s["meta"], 15, MUTED, False]]]}]
    return d

RENDERERS = {"cover": draw_cover, "toc": draw_toc, "section": draw_section,
             "cards": draw_cards, "compare": draw_compare, "strength": draw_strength,
             "quote": draw_quote, "chart": draw_chart, "flow": draw_flow,
             "goodbad": draw_goodbad, "columns3": draw_columns3,
             "summary": draw_summary, "qa": draw_qa, "thanks": draw_thanks}

# ============================================== 3a. SVG 后端（草稿 + 越界检查）
def esc(t):
    return (t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
             .replace('"', "&quot;"))

def render_svg(draw, page):
    lines = ['<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 960 540" font-family="sans-serif">']
    for it in draw:
        x, y, cx, cy = px(it["x"]), px(it["y"]), px(it["cx"]), px(it["cy"])
        k = it["k"]
        if k == "rect":
            lines.append(f'<rect x="{x}" y="{y}" width="{cx}" height="{cy}" fill="#{it.get("fill","FFFFFF")}"/>')
        elif k == "rrect":
            fill = f'#{it["fill"]}' if it.get("fill") else "none"
            stroke = (f'stroke="#{it["stroke"]}" stroke-width="{max(px(it.get("sw",12700)),0.7)}"'
                      if it.get("stroke") else "")
            lines.append(f'<rect x="{x}" y="{y}" width="{cx}" height="{cy}" rx="8" fill="{fill}" {stroke}/>')
        elif k == "oval":
            lines.append(f'<ellipse cx="{x + cx/2}" cy="{y + cy/2}" rx="{cx/2}" ry="{cy/2}" '
                         f'fill="none" stroke="#{it["stroke"]}" stroke-width="2"/>')
        elif k == "line":
            arrow = ' marker-end="url(#arr)"' if it.get("arrow") else ""
            if it.get("arrow") and not any("<defs>" in l for l in lines):
                lines.insert(1, '<defs><marker id="arr" markerWidth="8" markerHeight="8" '
                                'refX="6" refY="3" orient="auto"><path d="M0,0 L6,3 L0,6 z" '
                                f'fill="#{MUTED}"/></marker></defs>')
            lines.append(f'<line x1="{x}" y1="{y}" x2="{x + cx}" y2="{y + cy}" '
                         f'stroke="#{it["stroke"]}" stroke-width="2"{arrow}/>')
        elif k == "chart":
            lines.append(f'<rect x="{x}" y="{y}" width="{cx}" height="{cy}" rx="6" '
                         f'fill="#{SOFTBG}" stroke="#{LINE}"/>')
            lines.append(f'<text x="{x + 16}" y="{y + 28}" font-size="14" fill="#{DARK}">'
                         f'[chart] {esc(it["title"])}</text>')
        elif k == "text":
            align = it.get("align", "left")
            paras = it["runs"]
            sizes = [max(r[1] for r in p) for p in paras]
            th = sum(s * 1.5 for s in sizes)
            ty = y + (cy - th) / 2 if it.get("anchor") == "middle" else y
            for para in paras:
                h = max(r[1] for r in para) * 1.5
                anchor = {"left": "start", "center": "middle", "right": "end"}[align]
                tx2 = {"start": x + 4, "middle": x + cx / 2, "end": x + cx - 4}[anchor]
                for (t, sz, color, bold) in para:
                    fs = sz * 4 / 3
                    yy = ty + h * 0.8
                    for seg in t.split("\n"):
                        lines.append(f'<text x="{tx2:.0f}" y="{yy:.0f}" font-size="{fs:.0f}" '
                                     f'fill="#{color}" text-anchor="{anchor}" '
                                     f'font-weight="{"700" if bold else "400"}" '
                                     f'opacity="{it.get("alpha", 1.0)}">{esc(seg)}</text>')
                        yy += h
                ty += h
            if y < -4 or x < -4 or x + cx > 964 or y + cy > 544:
                lines.append(f'<!-- LINT: text block out of page at page {page} -->')
    lines.append("</svg>")
    return "\n".join(lines)

# ============================================== 3b. PPTX 后端
def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def set_font(run, size, color, bold):
    f = run.font
    f.name = FONT; f.size = Pt(size); f.color.rgb = _rgb(color); f.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

def render_pptx(prs, draw, blank):
    s = prs.slides.add_slide(blank)
    shapes = s.shapes
    for it in draw:
        k = it["k"]
        x, y, cx, cy = Emu(it["x"]), Emu(it["y"]), Emu(it["cx"]), Emu(it["cy"])
        if k == "rect":
            sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
            sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(it["fill"])
            sp.line.fill.background(); sp.shadow.inherit = False
        elif k == "rrect":
            sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
            try:
                sp.adjustments[0] = 0.06 if cy < 2000000 else 0.10
            except Exception:
                pass
            if it.get("fill"):
                sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(it["fill"])
            else:
                sp.fill.background()
            if it.get("stroke"):
                sp.line.color.rgb = _rgb(it["stroke"]); sp.line.width = Emu(it.get("sw", 12700))
            else:
                sp.line.fill.background()
            sp.shadow.inherit = False
        elif k == "oval":
            sp = shapes.add_shape(MSO_SHAPE.OVAL, x, y, cx, cy)
            sp.fill.background(); sp.line.color.rgb = _rgb(it["stroke"])
            sp.line.width = Emu(19050); sp.shadow.inherit = False
        elif k == "line":
            conn = shapes.add_connector(1, x, y, Emu(it["x"] + it["cx"]), Emu(it["y"] + it["cy"]))
            conn.line.color.rgb = _rgb(it["stroke"]); conn.line.width = Emu(19050)
            if it.get("arrow"):
                ln = conn.line._get_or_add_ln()
                ln.append(ln.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        elif k == "chart":
            data = CategoryChartData()
            data.categories = it["cats"]
            for name, vals in it["series"]:
                data.add_series(name, vals)
            gf = shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, data)
            gf.chart.has_title = True
            gf.chart.chart_title.text_frame.text = it["title"]
            gf.chart.has_legend = False
        elif k == "text":
            tb = shapes.add_textbox(x, y, cx, cy)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            if it.get("anchor") == "middle":
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for pi, para in enumerate(it["runs"]):
                p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                p.alignment = ALIGN[it.get("align", "left")]
                p.space_after = Pt(6)
                for (t, sz, color, bold) in para:
                    for j, seg in enumerate(t.split("\n")):
                        if j > 0:
                            p = tf.add_paragraph(); p.alignment = ALIGN[it.get("align", "left")]
                            p.space_after = Pt(6)
                        r = p.add_run(); r.text = seg; set_font(r, sz, color, bold)
    return s

# ============================================== main
def main():
    here = os.path.dirname(os.path.abspath(__file__))
    plan = plan_deck()
    # 阶段 1：版式规划落盘（写 pptx 前先审 plan）
    with open(os.path.join(here, "layout_plan.json"), "w") as f:
        json.dump(plan, f, ensure_ascii=False, indent=2)
    print(f"[plan ] {len(plan)} slides -> layout_plan.json")

    # 阶段 2：逐页 SVG 草稿
    svg_dir = os.path.join(here, "svg")
    os.makedirs(svg_dir, exist_ok=True)
    lint = []
    for spec in plan:
        svg = render_svg(RENDERERS[spec["layout"]](spec), spec["page"])
        with open(os.path.join(svg_dir, f"page-{spec['page']:02d}.svg"), "w") as f:
            f.write(svg)
        if "LINT" in svg:
            lint.append(spec["page"])
    print(f"[svg  ] {len(plan)} drafts -> svg/page-NN.svg" + (f"  (lint warn: {lint})" if lint else ""))

    # 阶段 3：渲染 pptx
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(W_EMU), Emu(H_EMU)
    blank = prs.slide_layouts[6]
    for spec in plan:
        render_pptx(prs, RENDERERS[spec["layout"]](spec), blank)
    out = os.path.join(here, "memory-nolinear-v2.pptx")
    prs.save(out)
    print(f"[pptx ] {out}  slides: {len(plan)}")

if __name__ == "__main__":
    main()
