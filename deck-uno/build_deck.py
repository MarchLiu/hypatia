#!/usr/bin/env python3
# Rebuild deck-uno/memory-nolinear.pptx with the visual language of deck/memory-nolinear-50min.pptx
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# palette (from reference deck)
NAVY   = RGBColor(0x16, 0x29, 0x4D)
BLUE   = RGBColor(0x2F, 0x6B, 0xFF)
TEAL   = RGBColor(0x12, 0xB5, 0xA5)
AMBER  = RGBColor(0xF2, 0xA9, 0x3B)
BG     = RGBColor(0xFA, 0xFB, 0xFD)
DARK   = RGBColor(0x1B, 0x2A, 0x41)
BODY   = RGBColor(0x44, 0x54, 0x6A)
MUTED  = RGBColor(0x7A, 0x8A, 0xA3)
LINE   = RGBColor(0xD9, 0xE1, 0xEF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0x9F, 0xB4, 0xD8)
ACCENTS = [BLUE, TEAL, AMBER]

W, H = Emu(12192000), Emu(6858000)
FONT = "Heiti SC"  # visible to fontconfig/headless LibreOffice on this host

def set_font(run, size, color, bold=False):
    f = run.font
    f.name = FONT; f.size = Pt(size); f.color.rgb = color; f.bold = bold
    r = run._r
    rPr = r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def no_fill(shape, line_color, w_emu=19050):
    shape.fill.background()
    shape.line.color.rgb = line_color; shape.line.width = Emu(w_emu)

def textbox(slide, x, y, cx, cy, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    """lines: list of (text, size, color, bold) or list-of-runs [(t,s,c,b),...]"""
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        runs = ln if isinstance(ln, list) else [ln]
        for (t, s, c, b) in runs:
            r = p.add_run(); r.text = t; set_font(r, s, c, b)
    return tb

def chip(slide, x, y, w=Emu(330200), h=Emu(50800), color=BLUE):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(sp, color)
    return sp

def deco_circles(slide):
    specs = [(7683500,1270000,4191000,BLUE),(8280400,1866900,2997200,TEAL),
             (8864600,2451100,1828800,AMBER)]
    for x,y,c,col in specs:
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x), Emu(y), Emu(c), Emu(c))
        no_fill(sp, col)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9728200), Emu(3314700), Emu(101600), Emu(101600))
    solid(dot, AMBER)

def bg_rect(slide, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    solid(sp, color)
    return sp

def card(slide, x, y, cx, cy, accent):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = LINE; sp.line.width = Emu(12700)
    sp.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Emu(0), y + Emu(60000), Emu(60000), cy - Emu(120000))
    solid(bar, accent)
    return sp

def split_lead(text):
    for sep in ("：", ":"):
        if sep in text:
            i = text.index(sep)
            return text[:i] + sep, text[i+1:]
    return None, text

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]

# ---------- content data (same 20 slides) ----------
cover = {
    "kicker": "HYPATIA · MEMORY DESIGN",
    "title": "非线性的记忆",
    "sub": "从线性对话流到可检索的知识结构",
    "meta": "基于 docs/memory-nolinear.md · 50 min",
}
slides = [
 ("为什么不能只保存对话", "开场", [
  "Context Length 的增长解决“读懂当前”，没有触碰“记住过去”",
  "保存对话本身 = 保存噪声：过程、寒暄、已被推翻的尝试",
  "正确目标：保存对话中被整理出来的知识",
  "线性记录三个痛点：体积膨胀、检索困难、无法关联"]),
 ("什么是一个 Session？", "第一步", [
  "粗看不成问题：一次启动到一次退出就是一段对话",
  "现实更复杂：一次会话可能换三个不相干的话题；一个任务可能横跨多次重启",
  "整理的单位若错了：后面所有摘要和抽取都对着错误的对象工作",
  "框架给出的物理边界 ≠ 语义上的 session"]),
 ("语义信号优先于物理边界", "Session 切分", [
  "话题切换（最强）：从 sqlite 存储跳到集群运维，就该切成两段",
  "任务边界（次强）：一次明确交付的完成 = 一段叙事的自然终结",
  "时间间隔（最弱）：中断一夜后的继续，值得检查是否切开",
  "物理边界仅作提示：收敛则顺势切分，未收敛建议并入下一段"]),
 ("切分粒度：一个 session，一个主题", "Session 切分", [
  "准则：一个 session 恰好容纳一个可独立陈述的主题",
  "太小：摘要被迫丢弃异质信息；太大：窗口没有空间等待话题收敛",
  "切分不追求一次定终身：低层原文总在，切错的边界可在更高层归档时重新缝合",
  "识别错误是可修复的，不识别才是不可修复的"]),
 ("对数空间的 Session Summary", "分层归档", [
  "若摘要和原对话一样长，摘要毫无意义",
  "目标：log₁₆(n) 的空间代价——每上升一层压缩 16 倍",
  "16 份对话 → 1 份会话摘要；16 份摘要 → 1 份主题摘要",
  "被压缩的是叙事和过程，不是结论"]),
 ("分层压缩 vs 一次性压缩", "分层归档", [
  "拒绝一次性压缩：“读完一万字写出两百字”损失不可控",
  "分层归档：每层只面对上一层的产物，输入已是高密度内容",
  "损失发生在每一层的边缘，而非一次性的坍塌",
  "B 树逻辑：不指望根节点存住所有数据，只指望从根出发总能找到"]),
 ("三层归档与下钻回溯", "分层归档", [
  "L0 原始对话：冷存，不删除，退出热路径",
  "L1 会话摘要：16:1 归档",
  "L2 主题摘要：再 16:1",
  "精确回溯：从高层逐层下钻，每下一层多一次检索"]),
 ("滑动窗口：抽取的时机", "知识抽取", [
  "对话进行中，hypatia-memory 用滑动窗口跟踪历史",
  "窗口内：待整理的候选；窗口滑过：已归档的既成事实",
  "知识点不是每句话都产生——十轮讨论才沉淀一个结论",
  "窗口的意义在于节奏，而非容量"]),
 ("窗口大小 = 话题的生命周期", "知识抽取", [
  "太短：把猜测和已被推翻的尝试误当成知识写入",
  "太长：整理的时机姗姗来迟",
  "实践：窗口覆盖“一个完整话题的自然生命周期”",
  "梳理时机：出现收敛迹象（结论被确认、话题被切换）"]),
 ("知识点：脱离上下文仍成立", "知识抽取", [
  "定义：一条自含的、脱离对话上下文仍然成立的陈述",
  "✗ “用户说要用 sqlite 换掉 duckdb”——那是聊天记录",
  "✓ “存储层评估过 sqlite 单库方案，因 json contains 查询不可行而搁置”",
  "前者只在语境里有意义，后者任何时候可检索、可直接使用"]),
 ("梳理时的裁决：新增、修正还是否定？", "知识抽取", [
  "每次梳理都要裁决：这条知识是新的，还是对既有知识的修正或否定？",
  "修正 → 更新既有条目",
  "否定 → 留下冲突的痕迹，绝不悄悄覆盖",
  "记忆系统最忌讳：悄悄忘记自己曾经知道什么"]),
 ("提取原则（一）：记什么", "提取原则", [
  "记录教训而非流水账：“用 Arc<Mutex<T>>”合格，“用合适的同步机制”不合格",
  "对质量保守、对时机激进：拿不准宁可跳过，但检查频率不能省",
  "纠错链是金子：问题→错答→纠正→修正，保留 initial attempt/错因/正确做法/lesson",
  "分类决定模板：一次答对、纠错链、探索讨论、bug 修复、设计决策各有骨架"]),
 ("提取原则（二）：怎么记", "提取原则", [
  "相对时间必须改写为绝对时间：“明天提醒”→“2026-08-29 提醒”",
  "知识点写给未来的读取者：“明天”这种词在入库那一刻就开始失效",
  "宁冗勿漏关系：哪怕谓词不完美，也要把三元组挂进图里",
  "绝不存储敏感信息：密码、密钥、token 不进入任何一层"]),
 ("三元组：把隐含关系显式化", "知识图", [
  "关系藏在语气和语序中：“sqlite 方案搁置”与“json contains 是核心依赖”有因果",
  "整理时必须显式化：每个知识点拆解为若干 (主语, 谓语, 宾语) 陈述，写入 RDF 风格的图",
  "分工：正文负责“读到时理解”，三元组负责“需要时找到”",
  "图的遍历是发现隐性关联的唯一手段"]),
 ("图遍历：发现隐性关联的唯一手段", "知识图", [
  "从 “sqlite” 出发沿关系走两步，撞见三年前关于 PG jsonb 的讨论",
  "这种相遇在任何全文检索里都不会发生",
  "承认噪声：三元组抽取不完美，谓词粒度难以精确",
  "一条错误的三元组是可发现、可挑战、可 resolve 的对象；未结构化的对话连犯错的资格都没有"]),
 ("自动回忆：按 Agent 机制适配", "读取通道", [
  "每轮重建 system prompt 的 Agent → 确定性注入：任务开始执行一次 JSE 查询",
  "有持续上下文、可自主调工具的 Agent → 回忆保持在技能层，需要时自己检索、中途补查",
  "被动响应、几乎无主动权的 Agent → 在对话边界（新任务、用户提醒）触发外部注入",
  "hypatia skill 的使命：同一套查询能力，包装成每个 Agent 最顺手的形态"]),
 ("显式写入与显式读取", "读取通道", [
  "显式写入：“把这份运维文档写进 hypatia”——现成材料跳过窗口等待，直接拆解入库；适合长时 Agent",
  "显式读取：“读一下 hypatia 找运维记忆”——检索从 Agent 的判断变成用户的指令",
  "显式读取是自动回忆失灵时的手动复位：成本一句话，回报是整个知识库的重新在场",
  "例：opencode 要直接 ssh 服务器时，一句提醒就够了"]),
 ("全景图：线性到非线性的完整路径", "总结", [
  "写入路径：线性对话流 → 语义切分 → L0/L1/L2 对数归档",
  "抽取路径：对话流经滑动窗口 → 知识点 + 三元组 → 知识图",
  "例外路径：显式写入直通知识图，跳过窗口",
  "读取通道：自动回忆（按 Agent 适配）+ 显式读取 → 服务 Agent 完成任务"]),
]

def add_cover():
    s = prs.slides.add_slide(blank)
    bg_rect(s, NAVY)
    deco_circles(s)
    chip(s, Emu(812800), Emu(711200))
    textbox(s, Emu(1300000), Emu(635000), Emu(6000000), Emu(300000),
            [(cover["kicker"], 15, SOFT, True)])
    textbox(s, Emu(812800), Emu(2286000), Emu(6800000), Emu(1800000),
            [(cover["title"], 64, WHITE, True)])
    textbox(s, Emu(812800), Emu(4000500), Emu(6400000), Emu(600000),
            [(cover["sub"], 24, SOFT, False)])
    textbox(s, Emu(812800), Emu(5900000), Emu(6400000), Emu(350000),
            [(cover["meta"], 15, MUTED, False)])

def add_toc():
    s = prs.slides.add_slide(blank)
    bg_rect(s, BG)
    header(s, "今天的路线", "OVERVIEW", 2)
    items = [
        ("01", "Session 的识别与切分", "语义信号：话题 · 任务 · 时间"),
        ("02", "对数空间的 Session Summary", "log₁₆(n) · 分层归档 · 下钻回溯"),
        ("03", "滑动窗口与知识点梳理", "时机 · 提取原则 · 纠错链"),
        ("04", "三元组与图结构", "显式化关系 · 图遍历发现"),
        ("05", "读写通道", "自动回忆适配 · 显式写入/读取"),
    ]
    top, bottom = Emu(1905000), Emu(6400000)
    ch = Emu(790000); gap = Emu(120000)
    y = top
    for i, (num, t, sub) in enumerate(items):
        acc = ACCENTS[i % 3]
        card(s, Emu(812800), y, Emu(10566400), ch, acc)
        textbox(s, Emu(1117600), y + Emu(180000), Emu(900000), Emu(450000),
                [(num, 28, acc, True)])
        textbox(s, Emu(2100000), y + Emu(110000), Emu(8800000), Emu(350000),
                [(t, 21, DARK, True)])
        textbox(s, Emu(2100000), y + Emu(450000), Emu(8800000), Emu(300000),
                [(sub, 15, MUTED, False)])
        y += ch + gap

def header(s, title, kicker, page_no):
    chip(s, Emu(812800), Emu(660400), color=BLUE)
    textbox(s, Emu(1300000), Emu(600000), Emu(6000000), Emu(280000),
            [(kicker, 13, BLUE, True)])
    textbox(s, Emu(812800), Emu(1000000), Emu(9500000), Emu(700000),
            [(title, 32, DARK, True)])
    textbox(s, Emu(10591800), Emu(660400), Emu(787400), Emu(280000),
            [(f"{page_no:02d} / 21", 13, MUTED, False)], align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(812800), Emu(1780000), Emu(10566400), Emu(12700))
    solid(ln, LINE)

def add_content(idx, title, kicker, bullets):
    s = prs.slides.add_slide(blank)
    bg_rect(s, BG)
    header(s, title, kicker, idx + 1)
    n = len(bullets)
    top, bottom = Emu(2032000), Emu(6550000)
    avail = bottom - top
    gap = Emu(150000)
    ch = Emu(int((avail - gap * (n - 1)) / n))
    size = 20 if n <= 4 else 18
    y = top
    for i, b in enumerate(bullets):
        acc = ACCENTS[i % 3]
        card(s, Emu(812800), y, Emu(10566400), ch, acc)
        lead, rest = split_lead(b)
        runs = []
        if lead:
            runs.append((lead, size, DARK, True))
        if rest:
            runs.append((rest, size, BODY, False))
        if not lead:
            runs = [(b, size, BODY, False)]
        tb = textbox(s, Emu(1231900), y + Emu(120000), Emu(9850000), ch - Emu(240000),
                     [runs], anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
        y += ch + gap

def add_closing():
    s = prs.slides.add_slide(blank)
    bg_rect(s, NAVY)
    deco_circles(s)
    chip(s, Emu(812800), Emu(711200), color=AMBER)
    textbox(s, Emu(1300000), Emu(635000), Emu(6000000), Emu(300000),
            [("写在最后", 15, SOFT, True)])
    textbox(s, Emu(812800), Emu(1900000), Emu(7000000), Emu(900000),
            [("对话是给人看的，知识是给机器用的", 40, WHITE, True)])
    lines = [
        ("对数分层", "控制体积"), ("滑动窗口", "控制时机"),
        ("三元组图", "提供发现路径"), ("双通道", "覆盖日常与例外"),
    ]
    for i, (a, b) in enumerate(lines):
        x = Emu(812800 + (i % 2) * 2921000)
        y = Emu(3150000 + (i // 2) * 950000)
        slide_pill(s, x, y, a, b)
    textbox(s, Emu(812800), Emu(5250000), Emu(9000000), Emu(500000),
            [("四件事拼在一起 = “有限记忆的 AI 使用无限长度的知识内容”", 20, SOFT, False)])
    textbox(s, Emu(812800), Emu(5950000), Emu(4000000), Emu(400000),
            [("Q & A", 28, WHITE, True)])

def slide_pill(s, x, y, a, b):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(2700000), Emu(850000))
    try:
        sp.adjustments[0] = 0.25
    except Exception:
        pass
    no_fill(sp, BLUE, 19050)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(120000)
    tf.margin_top = tf.margin_bottom = Emu(60000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = a; set_font(r, 18, WHITE, True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = b; set_font(r2, 13, SOFT, False)

add_cover()
add_toc()
for i, (t, k, bs) in enumerate(slides):
    add_content(i, t, k, bs)
add_closing()

out = "/Users/mars/jobs/hypatia/deck-uno/memory-nolinear.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
